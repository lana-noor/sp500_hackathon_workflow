# Copyright (c) Microsoft. All rights reserved.

"""
S&P 500 Quantitative Research Workflow — Five Sequential Agents
===============================================================

Pipeline:
  Agent 1  →  Query Planner          (FoundryChatClient)
               Converts the user question into a structured analytical plan:
               intent, scope, time window, analytical questions (for the CSV data),
               and qualitative questions (for web search).

  Agent 2  →  Quantitative Analyst   (FoundryAgent + Code Interpreter)
               Executes every analytical question from the plan against three CSV
               files (tickerlist, prices, returns). Returns structured findings
               with real computed numbers. No charts produced.

  Agent 3  →  Search Query Expander  (FoundryChatClient)
               Takes the original user query + Agent 2 findings summary and
               expands them into 5-8 precise web search query strings.

  Agent 4  →  Web Researcher         (FoundryChatClient + Tavily, per-query loop)
               Calls the Tavily Search API once per query from Agent 3, then uses
               FoundryChatClient to extract and structure evidence from the results.

  Agent 5  →  Synthesizer            (FoundryChatClient)
               Merges quantitative findings (Agent 2) and evidence pack (Agent 4)
               into a slide_deck_spec JSON. Every claim cites a number from Agent 2
               or a URL from Agent 4.

  Function →  generate_pptx()        (python-pptx — not an LLM)
               Reads slide_deck_spec from Agent 5 and builds the .pptx file.

Data files (in sp500_sample_data/):
  tickerlist.csv  — 505 S&P 500 constituents: name, ticker, sector
  returns.csv     — 60 monthly returns, Jan 2017 – Dec 2021, one column per ticker
  prices.csv      — ~62 monthly price snapshots, one column per ticker

Prerequisites:
  - QUANT_AGENT_NAME deployed in Azure AI Foundry with Code Interpreter enabled
    and the three CSV files attached to the agent in the Foundry portal
  - TAVILY_API_KEY set in .env (get a free key at tavily.com)
  - Azure credentials configured (DefaultAzureCredential)
  - Environment variables set (see .envsample)

Run:
  python portfolio_risk_workflow.py
  python portfolio_risk_workflow.py --query "How did tech compare to energy from 2017 to 2021?"
"""

import argparse
import asyncio
import json
import re
from datetime import datetime
from pathlib import Path
from typing import Optional

import os
from dotenv import load_dotenv

from agent_framework import Agent
from agent_framework.foundry import FoundryChatClient, FoundryAgent
from azure.identity import DefaultAzureCredential

# ---------------------------------------------------------------------------
# Optional: python-pptx for presentation generation
# ---------------------------------------------------------------------------
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("[Warning] python-pptx not installed — PPTX output will be skipped.")
    print("          Install with: pip install python-pptx")

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------
load_dotenv()

OUTPUT_DIR = Path(os.getenv("OUTPUT_DIR", "./output"))
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

SP500_DATA_DIR = Path(os.getenv("SP500_DATA_DIR", "./sp500_sample_data"))
PROMPTS_DIR    = Path(__file__).parent / "prompts"


def load_prompt(filename: str) -> str:
    path = PROMPTS_DIR / filename
    with open(path, "r", encoding="utf-8") as f:
        return f.read().strip()


# Agent references (deployed in Azure AI Foundry portal)
QUANT_AGENT         = os.getenv("QUANT_AGENT_NAME", "SP500QuantAgent")
QUANT_VERSION       = os.getenv("QUANT_AGENT_VERSION", "1")

TAVILY_API_KEY      = os.getenv("TAVILY_API_KEY", "")

# ---------------------------------------------------------------------------
# Load prompts
# ---------------------------------------------------------------------------
QUERY_PLANNER_INSTRUCTIONS        = load_prompt("query_planner_agent.txt")
SEARCH_EXPANDER_INSTRUCTIONS      = load_prompt("search_query_expander_agent.txt")
SYNTHESIZER_INSTRUCTIONS          = load_prompt("synthesizer_agent.txt")
WEB_RESEARCHER_INSTRUCTIONS       = load_prompt("web_researcher_agent.txt")
# Agent 2 prompt is set in the Foundry portal; kept here as a setup reference.
_QUANT_ANALYST_SYSTEM_PROMPT_REF  = load_prompt("quantitative_analyst_agent.txt")


# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
_DARK_NAVY   = (0x0D, 0x1B, 0x3E)
_NAVY        = (0x1A, 0x37, 0x6C)
_MED_BLUE    = (0x27, 0x5E, 0xAD)
_LIGHT_BLUE  = (0xD6, 0xE4, 0xF7)
_GOLD        = (0xF0, 0xA5, 0x00)
_GREEN       = (0x1E, 0x7E, 0x45)
_RED         = (0xC0, 0x39, 0x2B)
_WHITE       = (0xFF, 0xFF, 0xFF)
_LIGHT_GRAY  = (0xF4, 0xF6, 0xFA)
_MID_GRAY    = (0xC8, 0xD0, 0xDC)
_DARK_GRAY   = (0x30, 0x30, 0x30)

# ---------------------------------------------------------------------------
# Presentation builder — deterministic Python function (not an LLM)
# ---------------------------------------------------------------------------

def generate_pptx(slide_deck_spec: dict, output_path: Path) -> Path:
    """
    Build a PowerPoint presentation from slide_deck_spec produced by Agent 5.

    Args:
        slide_deck_spec: Structured JSON from Agent 5 (synthesizer).
        output_path:     Where to save the .pptx file.

    Returns:
        output_path after saving.
    """
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")

    def rgb(t):
        return RGBColor(*t)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    date_str = slide_deck_spec.get("analysis_date", datetime.now().strftime("%Y-%m-%d"))
    period   = slide_deck_spec.get("analysis_period", "2017–2021")
    footer_text = f"S&P 500 Research  |  {period}  |  {date_str}  |  Data: Jan 2017 – Dec 2021"

    # ------------------------------------------------------------------
    # Shared drawing helpers
    # ------------------------------------------------------------------

    def add_rect(slide, l, t, w, h, fill):
        sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = rgb(fill)
        sh.line.fill.background()
        return sh

    def add_text(slide, text, l, t, w, h, size=14, bold=False,
                 color=_WHITE, align=PP_ALIGN.LEFT, wrap=True):
        txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        txb.text_frame.word_wrap = wrap
        p = txb.text_frame.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        return txb

    def add_header(slide, title, subtitle=""):
        add_rect(slide, 0, 0, 13.33, 1.25, _NAVY)
        add_rect(slide, 0, 1.25, 13.33, 0.07, _GOLD)
        add_text(slide, title, 0.35, 0.12, 11.5, 0.85, size=26, bold=True)
        if subtitle:
            add_text(slide, subtitle, 0.35, 0.88, 11.5, 0.38, size=12, color=_LIGHT_BLUE)

    def add_footer(slide):
        add_rect(slide, 0, 7.12, 13.33, 0.38, _LIGHT_GRAY)
        add_text(slide, footer_text, 0.3, 7.14, 12.5, 0.35,
                 size=8, color=_DARK_GRAY, align=PP_ALIGN.LEFT)

    def add_bullets(slide, items, l, t, w, h, size=13, color=_DARK_GRAY):
        txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = txb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(5)
            run = p.add_run()
            run.text = f"•   {item}"
            run.font.size = Pt(size)
            run.font.color.rgb = rgb(color)

    def add_table_slide(slide, columns, rows, l, t, w, h):
        if not columns or not rows:
            return
        col_count = len(columns)
        row_count = len(rows) + 1
        tbl = slide.shapes.add_table(
            row_count, col_count, Inches(l), Inches(t), Inches(w), Inches(h)
        ).table
        col_w = Inches(w / col_count)
        for c in tbl.columns:
            c.width = col_w
        # Header
        for ci, name in enumerate(columns):
            cell = tbl.cell(0, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(_NAVY)
            cell.text = str(name)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.bold = True
            p.runs[0].font.size = Pt(11)
            p.runs[0].font.color.rgb = rgb(_WHITE)
        # Data rows
        for ri, row in enumerate(rows):
            fill = _LIGHT_BLUE if ri % 2 == 0 else _WHITE
            for ci in range(col_count):
                cell = tbl.cell(ri + 1, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(fill)
                val = str(row[ci]) if ci < len(row) else ""
                cell.text = val
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.runs[0].font.size = Pt(10)
                p.runs[0].font.color.rgb = rgb(_DARK_GRAY)

    # ------------------------------------------------------------------
    # Render each slide from the spec
    # ------------------------------------------------------------------
    slides_spec = slide_deck_spec.get("slides", [])

    for slide_spec in slides_spec:
        stype = slide_spec.get("slide_type", "narrative")
        s = prs.slides.add_slide(blank)

        # ── Title slide ──────────────────────────────────────────────
        if stype == "title":
            add_rect(s, 0, 0, 13.33, 7.5, _DARK_NAVY)
            add_rect(s, 0, 3.3, 13.33, 0.09, _GOLD)
            title = slide_spec.get("title", slide_deck_spec.get("presentation_title", ""))
            subtitle = slide_spec.get("subtitle", "")
            meta = slide_spec.get("meta", "")
            add_text(s, title, 0.8, 0.9, 11.7, 2.1, size=36, bold=True,
                     align=PP_ALIGN.CENTER)
            if subtitle:
                add_text(s, subtitle, 0.8, 3.6, 11.7, 1.0, size=18,
                         color=_LIGHT_BLUE, align=PP_ALIGN.CENTER)
            if meta:
                add_text(s, meta, 0.8, 4.8, 11.7, 0.6, size=13,
                         color=(0x90, 0xA8, 0xC8), align=PP_ALIGN.CENTER)
            add_text(s, "Source data: S&P 500 constituents Jan 2017 – Dec 2021",
                     0.8, 6.6, 11.7, 0.5, size=10,
                     color=(0x60, 0x78, 0x98), align=PP_ALIGN.CENTER)

        # ── Executive summary ────────────────────────────────────────
        elif stype == "exec_summary":
            add_rect(s, 0, 0, 13.33, 7.5, _LIGHT_GRAY)
            add_header(s, slide_spec.get("title", "Executive Summary"))
            bullets = slide_spec.get("bullets", [])
            # Show bullets as coloured cards
            card_colors = [_NAVY, _MED_BLUE, (0x35, 0x6B, 0xA8),
                           (0x1E, 0x5A, 0x9C), (0x16, 0x4D, 0x8C)]
            for i, bullet in enumerate(bullets[:5]):
                ly = 1.5 + i * 1.05
                add_rect(s, 0.4, ly, 12.5, 0.88, card_colors[i % len(card_colors)])
                add_text(s, bullet, 0.7, ly + 0.12, 12.0, 0.72, size=14)

            add_footer(s)

        # ── Comparison table ─────────────────────────────────────────
        elif stype == "comparison_table":
            add_rect(s, 0, 0, 13.33, 7.5, _LIGHT_GRAY)
            add_header(s, slide_spec.get("title", ""))
            tbl_data = slide_spec.get("table", {})
            columns = tbl_data.get("columns", [])
            rows    = tbl_data.get("rows", [])
            note    = slide_spec.get("note", "")
            table_h = min(5.0, 0.55 + len(rows) * 0.48)
            add_table_slide(s, columns, rows, 0.4, 1.45, 12.5, table_h)
            if note:
                add_text(s, note, 0.4, 1.45 + table_h + 0.1, 12.5, 0.45,
                         size=10, color=_DARK_GRAY)
            add_footer(s)

        # ── Narrative ────────────────────────────────────────────────
        elif stype == "narrative":
            add_rect(s, 0, 0, 13.33, 7.5, _LIGHT_GRAY)
            add_header(s, slide_spec.get("title", ""))
            bullets   = slide_spec.get("bullets", [])
            citations = slide_spec.get("citations", [])
            add_rect(s, 0.4, 1.44, 12.5, 5.0, _WHITE)
            add_bullets(s, bullets, 0.7, 1.6, 12.0, 4.8, size=14)
            if citations:
                cite_str = "Sources: " + " | ".join(citations[:4])
                add_text(s, cite_str, 0.4, 6.55, 12.5, 0.45, size=8,
                         color=(0x60, 0x78, 0x98))
            add_footer(s)

        # ── Sources slide ────────────────────────────────────────────
        elif stype == "sources":
            add_rect(s, 0, 0, 13.33, 7.5, _DARK_NAVY)
            add_rect(s, 0, 0, 13.33, 1.25, (0x0A, 0x12, 0x28))
            add_rect(s, 0, 1.25, 13.33, 0.07, _GOLD)
            add_text(s, slide_spec.get("title", "Sources & References"),
                     0.35, 0.15, 12.5, 0.95, size=26, bold=True)
            sources = slide_spec.get("sources", [])
            for i, src in enumerate(sources[:9]):
                ly = 1.5 + i * 0.58
                add_rect(s, 0.4, ly, 12.5, 0.48, (0x18, 0x2E, 0x56))
                label = f"{src.get('title', 'Source')}  ({src.get('date', '')})"
                add_text(s, label, 0.6, ly + 0.04, 6.5, 0.38, size=11, color=_LIGHT_BLUE)
                add_text(s, src.get("url", ""), 7.2, ly + 0.04, 5.6, 0.38,
                         size=9, color=(0x90, 0xA8, 0xC8))
            add_text(s, "Data source: S&P 500 constituent data Jan 2017 – Dec 2021  |  For internal use only.",
                     0.4, 7.0, 12.5, 0.4, size=8, color=(0x60, 0x78, 0x98))

        else:
            # Fallback: treat as narrative
            add_rect(s, 0, 0, 13.33, 7.5, _LIGHT_GRAY)
            add_header(s, slide_spec.get("title", stype))
            bullets = slide_spec.get("bullets", slide_spec.get("body_bullets", []))
            add_bullets(s, bullets, 0.4, 1.5, 12.5, 5.5, size=13)
            add_footer(s)

    prs.save(str(output_path))
    return output_path


# ---------------------------------------------------------------------------
# Checkpoint store — single JSON file, one key per completed step
# ---------------------------------------------------------------------------

class CheckpointStore:
    """
    Persists workflow progress to a single JSON file after each agent completes.

    File layout:
    {
      "query":          "<original user query>",
      "started_at":     "<ISO timestamp>",
      "last_updated":   "<ISO timestamp>",
      "completed_steps": ["agent1", "agent2", ...],
      "agent1_plan":        "<raw text output from Agent 1>",
      "agent2_findings":    "<raw text output from Agent 2>",
      "agent3_queries":     "<raw text output from Agent 3>",
      "agent4_evidence":    "<raw text output from Agent 4>",
      "agent5_deck_spec":   "<raw text output from Agent 5>",
      "pptx_path":          "<path or null>"
    }

    If the file already exists and contains a matching query, execution resumes
    from the first incomplete step — completed steps are skipped entirely.
    """

    STEPS = ["agent1", "agent2", "agent3", "agent4", "agent5"]
    KEY_MAP = {
        "agent1": "agent1_plan",
        "agent2": "agent2_findings",
        "agent3": "agent3_queries",
        "agent4": "agent4_evidence",
        "agent5": "agent5_deck_spec",
    }

    def __init__(self, path: Path, query: str) -> None:
        self.path  = path
        self.query = query
        self._data: dict = {}
        self._load_or_init()

    def _load_or_init(self) -> None:
        if self.path.exists():
            try:
                self._data = json.loads(self.path.read_text(encoding="utf-8"))
                saved_query = self._data.get("query", "")
                if saved_query != self.query:
                    print(f"[Checkpoint] Query mismatch — starting fresh.")
                    print(f"  Saved:   {saved_query[:80]}")
                    print(f"  Current: {self.query[:80]}")
                    self._data = {}
                else:
                    done = self._data.get("completed_steps", [])
                    print(f"[Checkpoint] Loaded existing run. Completed steps: {done}")
            except (json.JSONDecodeError, OSError) as exc:
                print(f"[Checkpoint] Could not load checkpoint ({exc}) — starting fresh.")
                self._data = {}

        if not self._data:
            self._data = {
                "query":           self.query,
                "started_at":      datetime.now().isoformat(),
                "completed_steps": [],
            }
            self._flush()

    def is_done(self, step: str) -> bool:
        return step in self._data.get("completed_steps", [])

    def get(self, step: str) -> Optional[str]:
        return self._data.get(self.KEY_MAP[step])

    def save(self, step: str, text: str) -> None:
        self._data[self.KEY_MAP[step]] = text
        completed = self._data.setdefault("completed_steps", [])
        if step not in completed:
            completed.append(step)
        self._data["last_updated"] = datetime.now().isoformat()
        self._flush()
        print(f"[Checkpoint] Step '{step}' saved -> {self.path}")

    def save_outputs(self, pptx_path: Optional[Path]) -> None:
        self._data["pptx_path"]    = str(pptx_path) if pptx_path else None
        self._data["last_updated"] = datetime.now().isoformat()
        self._flush()

    def _flush(self) -> None:
        self.path.write_text(
            json.dumps(self._data, indent=2, ensure_ascii=False),
            encoding="utf-8",
        )

    @property
    def all_data(self) -> dict:
        return self._data


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _extract_text_from_result(result) -> str:
    if hasattr(result, "text") and result.text:
        return result.text
    try:
        if hasattr(result, "messages") and result.messages:
            for message in reversed(result.messages):
                if hasattr(message, "contents"):
                    for content in message.contents:
                        if getattr(content, "type", None) == "text":
                            if hasattr(content, "text") and content.text:
                                return content.text
    except Exception:
        pass
    return str(result)


def _banner(title: str) -> None:
    bar = "=" * 65
    print(f"\n{bar}")
    print(f"  {title}")
    print(bar)


def _parse_json(text: str) -> Optional[dict]:
    """Extract and parse the first JSON object found in a text string."""
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass
    match = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", text, re.DOTALL)
    if match:
        try:
            return json.loads(match.group(1))
        except json.JSONDecodeError:
            pass
    match = re.search(r"\{.*\}", text, re.DOTALL)
    if match:
        try:
            return json.loads(match.group(0))
        except json.JSONDecodeError:
            pass
    return None


# ---------------------------------------------------------------------------
# Agent 2 output validation
# ---------------------------------------------------------------------------

#: Strings that indicate Agent 2 returned a template instead of real data.
_HALLUCINATION_MARKERS = [
    "calculated value",
    "[number_of_rows]",
    "[number_of_columns]",
    "[start_date]",
    "[end_date]",
    "[value]",
    "[ticker]",
    "[sector]",
    "placeholder",
    "your_value",
    "<value>",
    "<number>",
    "INSERT",
]


def _validate_agent2_output(text: str) -> None:
    """
    Raise ValueError if Agent 2's output looks like a hallucinated template
    rather than real computed results.

    Checks:
      1. Output is valid JSON.
      2. Required key 'findings' is present and non-empty.
      3. No hallucination marker strings appear in the text.
      4. findings list contains at least one entry with a real answer.
    """
    parsed = _parse_json(text)
    if parsed is None:
        raise ValueError(
            "[Agent 2 Validation] Output is not valid JSON. "
            "Agent may have returned prose instead of JSON. "
            "Check that the Code Interpreter tool ran successfully."
        )

    # ── Check for hallucination marker strings ───────────────────────
    text_lower = text.lower()
    for marker in _HALLUCINATION_MARKERS:
        if marker.lower() in text_lower:
            raise ValueError(
                f"[Agent 2 Validation] Output contains placeholder text: '{marker}'. "
                "Agent did not execute code — it returned a template. "
                "Ensure the Code Interpreter tool is enabled and the CSV files are attached."
            )

    # ── findings must be present ─────────────────────────────────────
    if "findings" not in parsed:
        raise ValueError(
            f"[Agent 2 Validation] Missing required key 'findings'. "
            f"Present keys: {list(parsed.keys())}."
        )

    findings = parsed.get("findings")

    # findings can be a list (new schema) or dict (old schema) — handle both
    if isinstance(findings, list):
        if len(findings) == 0:
            raise ValueError(
                "[Agent 2 Validation] 'findings' list is empty — no questions were answered."
            )
        # Check at least one finding has a real answer
        answered = [f for f in findings if isinstance(f, dict) and f.get("answer", "").strip()]
        if not answered:
            raise ValueError(
                "[Agent 2 Validation] All findings have empty answers — agent did not compute results."
            )
    elif isinstance(findings, dict):
        if not findings:
            raise ValueError(
                "[Agent 2 Validation] 'findings' dict is empty — no results were computed."
            )
        non_empty = [v for v in findings.values() if v not in (None, "", [], {})]
        if not non_empty:
            raise ValueError(
                "[Agent 2 Validation] All values in 'findings' are empty or null."
            )
    else:
        raise ValueError(
            f"[Agent 2 Validation] 'findings' must be a list or dict, got: {type(findings)}"
        )

    n = len(findings) if isinstance(findings, (list, dict)) else 0
    print(f"[Agent 2 Validation] OK — {n} finding(s) validated.")


# ---------------------------------------------------------------------------
# Workflow Executor
# ---------------------------------------------------------------------------

class SP500AnalysisWorkflowExecutor:
    """
    Sequential 5-agent workflow for S&P 500 quantitative research.

    Step 1  — Query Planner (FoundryChatClient): interprets user query into research plan.
    Step 2  — Quantitative Analyst (FoundryAgent + Code Interpreter): answers quantitative
              questions using tickerlist.csv, prices.csv, returns.csv.
    Step 3  — Search Query Expander (FoundryChatClient): converts qualitative questions
              and data findings into precise web search queries.
    Step 4  — Web Researcher (FoundryAgent + Bing/Tavily): executes searches, returns
              evidence pack with URLs.
    Step 5  — Synthesizer (FoundryChatClient): merges findings + evidence into
              slide_deck_spec JSON.
    Function— generate_pptx(): builds the final .pptx (python-pptx).
    """

    def __init__(self) -> None:
        self.credential       = DefaultAzureCredential()
        self.project_endpoint = os.getenv("AZURE_AI_PROJECT_ENDPOINT")
        self.model            = os.getenv("AZURE_AI_MODEL_DEPLOYMENT_NAME")

    async def execute(self, user_query: str, run_id: Optional[str] = None) -> dict:
        """
        Run the full S&P 500 research workflow with checkpoint/resume support.

        All agent outputs are written into a single JSON file after each step.
        If the run is interrupted, re-running with the same query (and same run_id
        if provided) will skip already-completed steps and resume from the next one.

        Args:
            user_query: Natural language research question.
            run_id:     Optional stable identifier for this run. If omitted, a
                        timestamp is generated on first run and reused on resume.
                        Pass the same run_id to resume a specific interrupted run.

        Returns:
            dict with all agent outputs and final output paths.
        """
        # ── Resolve checkpoint file path ─────────────────────────────
        # The checkpoint file IS the output file — same single JSON for everything.
        if run_id is None:
            # Look for an existing checkpoint for this query first
            existing = sorted(OUTPUT_DIR.glob("sp500_analysis_*.json"), reverse=True)
            for candidate in existing:
                try:
                    data = json.loads(candidate.read_text(encoding="utf-8"))
                    if data.get("query") == user_query:
                        run_id = candidate.stem.replace("sp500_analysis_", "")
                        print(f"[Checkpoint] Found existing run: {candidate.name}")
                        break
                except (json.JSONDecodeError, OSError):
                    continue
            if run_id is None:
                run_id = datetime.now().strftime("%Y%m%d_%H%M%S")

        checkpoint_path = OUTPUT_DIR / f"sp500_analysis_{run_id}.json"
        ckpt = CheckpointStore(checkpoint_path, user_query)

        _banner("S&P 500 RESEARCH WORKFLOW — START (5 AGENTS)")
        print(f"  Query:      {user_query}")
        print(f"  Run ID:     {run_id}")
        print(f"  Checkpoint: {checkpoint_path}\n")

        # ── Agent 1: Query Planner ───────────────────────────────────
        if ckpt.is_done("agent1"):
            print("[Agent 1] SKIPPED (checkpoint)")
            plan_text = ckpt.get("agent1")
        else:
            plan_text = await self._run_query_planner(user_query)
            ckpt.save("agent1", plan_text)

        # ── Agent 2: Quantitative Analyst ────────────────────────────
        if ckpt.is_done("agent2"):
            print("[Agent 2] SKIPPED (checkpoint)")
            findings_text = ckpt.get("agent2")
        else:
            findings_text = await self._run_quantitative_analyst(plan_text)
            ckpt.save("agent2", findings_text)

        # ── Agent 3: Search Query Expander ───────────────────────────
        if ckpt.is_done("agent3"):
            print("[Agent 3] SKIPPED (checkpoint)")
            queries_text = ckpt.get("agent3")
        else:
            queries_text = await self._run_search_query_expander(
                user_query, plan_text, findings_text
            )
            ckpt.save("agent3", queries_text)

        # ── Agent 4: Web Researcher ──────────────────────────────────
        if ckpt.is_done("agent4"):
            print("[Agent 4] SKIPPED (checkpoint)")
            evidence_text = ckpt.get("agent4")
        else:
            evidence_text = await self._run_web_researcher(queries_text)
            ckpt.save("agent4", evidence_text)

        # ── Agent 5: Synthesizer ─────────────────────────────────────
        if ckpt.is_done("agent5"):
            print("[Agent 5] SKIPPED (checkpoint)")
            deck_spec_text = ckpt.get("agent5")
        else:
            deck_spec_text = await self._run_synthesizer(
                plan_text, findings_text, evidence_text
            )
            ckpt.save("agent5", deck_spec_text)

        # ── PPTX generation (local function) ─────────────────────────
        pptx_path: Optional[Path] = None
        deck_spec = _parse_json(deck_spec_text)

        if deck_spec and PPTX_AVAILABLE:
            pptx_out = OUTPUT_DIR / f"sp500_analysis_{run_id}.pptx"
            try:
                generate_pptx(deck_spec, pptx_out)
                pptx_path = pptx_out
                print(f"[Executor] PPTX saved -> {pptx_path}")
            except Exception as exc:
                print(f"[Executor] PPTX generation failed: {exc}")
        elif not PPTX_AVAILABLE:
            print("[Executor] Skipping PPTX — python-pptx not installed.")
        else:
            print("[Executor] Skipping PPTX — could not parse slide_deck_spec from Agent 5.")

        # ── Persist final output path into the checkpoint file ───────
        ckpt.save_outputs(pptx_path)

        _banner("WORKFLOW COMPLETE")
        print(f"  Output JSON -> {checkpoint_path}")
        if pptx_path:
            print(f"  Output PPTX -> {pptx_path}")

        return {
            "research_plan":   plan_text,
            "findings":        findings_text,
            "search_queries":  queries_text,
            "evidence_pack":   evidence_text,
            "slide_deck_spec": deck_spec_text,
            "pptx_path":       str(pptx_path) if pptx_path else None,
            "json_path":       str(checkpoint_path),
        }

    # ------------------------------------------------------------------
    # Agent 1 — Query Planner (FoundryChatClient)
    # ------------------------------------------------------------------

    async def _run_query_planner(self, user_query: str) -> str:
        print("\n[Agent 1] Query Planner (FoundryChatClient)...")
        print("[Agent 1] Converting user question into structured research plan...")

        agent = Agent(
            client=FoundryChatClient(
                project_endpoint=self.project_endpoint,
                model=self.model,
                credential=self.credential,
            ),
            instructions=QUERY_PLANNER_INSTRUCTIONS,
        )
        result = await agent.run(
            f"Produce a research plan for the following user query:\n\n{user_query}"
        )
        text = _extract_text_from_result(result)
        print(f"[Agent 1] Done. Research plan produced.\n{text[:200]}...\n")
        return text

    # ------------------------------------------------------------------
    # Agent 2 — Quantitative Analyst (FoundryAgent + Code Interpreter)
    # ------------------------------------------------------------------

    async def _run_quantitative_analyst(self, plan_text: str) -> str:
        print(f"\n[Agent 2] Quantitative Analyst — calling {QUANT_AGENT} v{QUANT_VERSION}...")
        print("[Agent 2] Running Python analysis against tickerlist, prices, returns CSVs...")

        agent = FoundryAgent(
            project_endpoint=self.project_endpoint,
            agent_name=QUANT_AGENT,
            agent_version=str(QUANT_VERSION),
            credential=self.credential,
        )
        # The CSV files (tickerlist.csv, prices.csv, returns.csv) must be
        # pre-attached to this agent in the Azure AI Foundry portal.
        user_message = (
            "Please run a quantitative analysis using the Python code interpreter "
            "and the attached CSV files. The research plan below specifies what to compute.\n\n"
            "=== RESEARCH PLAN ===\n"
            f"{plan_text}\n\n"
            "Instructions:\n"
            "1. Start by loading all three CSV files and printing their shapes and column names "
            "to confirm the data loaded correctly.\n"
            "2. Answer every question listed in analytical_questions using Python code and real "
            "computed values from the data. Do not estimate or approximate.\n"
            "3. Do not produce any charts or visualizations — data analysis only.\n"
            "4. If a question cannot be answered from the CSV files, add it to out_of_scope.\n"
            "5. Return a single valid JSON object with keys: "
            "analysis_date, data_coverage, findings, data_summary, out_of_scope.\n"
            "6. Return only the JSON object — no prose, no markdown fences around it."
        )
        result = await agent.run(user_message)
        text = _extract_text_from_result(result)

        # Validate the output — catch hallucinated/template responses early
        _validate_agent2_output(text)

        print(f"[Agent 2] Done. Quantitative analysis complete.\n{text[:200]}...\n")
        return text

    # ------------------------------------------------------------------
    # Agent 3 — Search Query Expander (FoundryChatClient)
    # ------------------------------------------------------------------

    async def _run_search_query_expander(
        self, user_query: str, plan_text: str, findings_text: str
    ) -> str:
        print("\n[Agent 3] Search Query Expander (FoundryChatClient)...")
        print("[Agent 3] Expanding user query into targeted web search queries...")

        agent = Agent(
            client=FoundryChatClient(
                project_endpoint=self.project_endpoint,
                model=self.model,
                credential=self.credential,
            ),
            instructions=SEARCH_EXPANDER_INSTRUCTIONS,
        )
        user_message = (
            "Generate web search queries for the following research request.\n\n"
            "=== ORIGINAL USER QUERY ===\n"
            f"{user_query}\n\n"
            "=== RESEARCH PLAN (Agent 1) ===\n"
            f"{plan_text}\n\n"
            "=== QUANTITATIVE FINDINGS SUMMARY (Agent 2) ===\n"
            f"{findings_text[:1500]}\n\n"
            "Return ONLY valid JSON."
        )
        result = await agent.run(user_message)
        text = _extract_text_from_result(result)
        print(f"[Agent 3] Done. Search queries generated.\n{text[:200]}...\n")
        return text

    # ------------------------------------------------------------------
    # Agent 4 — Web Researcher (FoundryChatClient + Tavily, one call per query)
    # ------------------------------------------------------------------

    async def _search_tavily(self, query: str) -> tuple[list[dict], str]:
        """
        Call the Tavily Search REST API for a single query.
        Returns (results_list, answer_string).
        Falls back to empty results if TAVILY_API_KEY is not set.
        """
        if not TAVILY_API_KEY:
            print(f"[Agent 4]   TAVILY_API_KEY not set — skipping: {query[:60]}")
            return [], ""
        try:
            import httpx
            async with httpx.AsyncClient(timeout=30) as client:
                resp = await client.post(
                    "https://api.tavily.com/search",
                    json={
                        "api_key": TAVILY_API_KEY,
                        "query": query,
                        "search_depth": "basic",
                        "max_results": 3,
                        "include_answer": True,
                    },
                )
                resp.raise_for_status()
                data = resp.json()
                return data.get("results", []), data.get("answer", "")
        except Exception as exc:
            print(f"[Agent 4]   Tavily search failed ({exc}) — skipping query.")
            return [], ""

    async def _run_web_researcher(self, queries_text: str) -> str:
        print("\n[Agent 4] Web Researcher (FoundryChatClient + Tavily, per-query loop)...")

        queries_json = _parse_json(queries_text) or {}
        queries = queries_json.get("queries", [])

        if not queries:
            print("[Agent 4] No queries found in Agent 3 output — skipping web research.")
            return json.dumps({"evidence_pack": [], "coverage_gaps": ["No queries received from Agent 3"]})

        evidence_pack: list[dict] = []
        coverage_gaps: list[str] = []

        for i, query_obj in enumerate(queries, 1):
            query_id  = query_obj.get("query_id", f"q{i}")
            query_str = query_obj.get("query_string", "")
            serves    = query_obj.get("serves_question", "")

            print(f"[Agent 4] Query {i}/{len(queries)}: {query_str[:80]}...")

            # Step 1 — fetch search results from Tavily
            results, tavily_answer = await self._search_tavily(query_str)

            if not results and not tavily_answer:
                coverage_gaps.append(serves or query_str)
                evidence_pack.append({
                    "query_id": query_id,
                    "query_string": query_str,
                    "serves_question": serves,
                    "sources": [],
                    "synthesis_note": "",
                    "coverage_gap": True,
                })
                continue

            # Step 2 — use FoundryChatClient to extract and structure evidence
            agent = Agent(
                client=FoundryChatClient(
                    project_endpoint=self.project_endpoint,
                    model=self.model,
                    credential=self.credential,
                ),
                instructions=WEB_RESEARCHER_INSTRUCTIONS,
            )

            search_context = (
                f"Query ID: {query_id}\n"
                f"Search query: {query_str}\n"
                f"Serves question: {serves}\n\n"
                f"Tavily answer: {tavily_answer}\n\n"
                f"Search results:\n"
                + json.dumps(results, indent=2)
            )

            try:
                result = await agent.run(search_context)
                text = _extract_text_from_result(result)
                evidence_obj = _parse_json(text)
                if evidence_obj:
                    evidence_obj.setdefault("query_id", query_id)
                    evidence_obj.setdefault("query_string", query_str)
                    evidence_obj.setdefault("serves_question", serves)
                    evidence_pack.append(evidence_obj)
                else:
                    # Could not parse — store raw
                    evidence_pack.append({
                        "query_id": query_id,
                        "query_string": query_str,
                        "serves_question": serves,
                        "sources": [],
                        "synthesis_note": text[:400],
                        "coverage_gap": False,
                    })
            except Exception as exc:
                print(f"[Agent 4]   Evidence extraction failed: {exc}")
                coverage_gaps.append(serves or query_str)

        print(f"[Agent 4] Done. {len(evidence_pack)} evidence entries, "
              f"{len(coverage_gaps)} coverage gap(s).\n")

        return json.dumps(
            {"evidence_pack": evidence_pack, "coverage_gaps": coverage_gaps},
            indent=2, ensure_ascii=False
        )

    # ------------------------------------------------------------------
    # Agent 5 — Synthesizer (FoundryChatClient)
    # ------------------------------------------------------------------

    async def _run_synthesizer(
        self,
        plan_text: str,
        findings_text: str,
        evidence_text: str,
    ) -> str:
        print("\n[Agent 5] Synthesizer (FoundryChatClient)...")
        print("[Agent 5] Merging findings + evidence into slide_deck_spec...")

        agent = Agent(
            client=FoundryChatClient(
                project_endpoint=self.project_endpoint,
                model=self.model,
                credential=self.credential,
            ),
            instructions=SYNTHESIZER_INSTRUCTIONS,
        )
        user_message = (
            "Build the slide_deck_spec JSON from the inputs below.\n\n"
            "=== RESEARCH PLAN (Agent 1) ===\n"
            f"{plan_text}\n\n"
            "=== QUANTITATIVE FINDINGS (Agent 2) ===\n"
            f"{findings_text}\n\n"
            "=== EVIDENCE PACK (Agent 4) ===\n"
            f"{evidence_text}\n\n"
            "Rules:\n"
            "- Every factual claim must cite a number from Agent 2's findings OR a URL from Agent 4.\n"
            "- Include 5-8 slides total (title + exec_summary + body slides + sources).\n"
            "- Return ONLY valid JSON matching your instructions schema."
        )
        result = await agent.run(user_message)
        text = _extract_text_from_result(result)
        print("[Agent 5] Done. Slide deck spec generated.\n")
        return text


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

DEFAULT_QUERY = (
    "How did Information Technology compare to Energy from 2017 to 2021, "
    "and what drove the differences in performance?"
)


def _delete_run(run_id: str | None, query: str | None) -> None:
    """
    Delete the checkpoint JSON (and associated PPTX) for a given run.

    If run_id is provided, delete that specific file.
    Otherwise find the latest checkpoint whose query matches.
    """
    targets: list[Path] = []

    if run_id:
        candidate = OUTPUT_DIR / f"sp500_analysis_{run_id}.json"
        if candidate.exists():
            targets.append(candidate)
        else:
            print(f"[Reset] No checkpoint found for run-id '{run_id}' ({candidate})")
            return
    else:
        existing = sorted(OUTPUT_DIR.glob("sp500_analysis_*.json"), reverse=True)
        for candidate in existing:
            try:
                data = json.loads(candidate.read_text(encoding="utf-8"))
                if query and data.get("query") == query:
                    targets.append(candidate)
                    break
                elif not query:
                    targets.append(candidate)
                    break
            except (json.JSONDecodeError, OSError):
                continue

    if not targets:
        print("[Reset] No matching checkpoint found — nothing to delete.")
        return

    for json_path in targets:
        stem = json_path.stem  # e.g. sp500_analysis_20260427_143022
        run_id_found = stem.replace("sp500_analysis_", "")

        # Delete JSON checkpoint
        json_path.unlink()
        print(f"[Reset] Deleted checkpoint: {json_path}")

        # Delete PPTX if it exists
        pptx_path = OUTPUT_DIR / f"sp500_analysis_{run_id_found}.pptx"
        if pptx_path.exists():
            pptx_path.unlink()
            print(f"[Reset] Deleted PPTX:        {pptx_path}")

        # Delete charts referenced in the checkpoint (already deleted JSON, so skip re-reading)
        print("[Reset] Done.")

    print("[Reset] Done. Run again without --reset to start fresh.")


async def main() -> None:
    parser = argparse.ArgumentParser(
        description="S&P 500 Quantitative Research Workflow — 5 Sequential Agents",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Fresh run with default query
  python portfolio_risk_workflow.py

  # Custom query
  python portfolio_risk_workflow.py --query "How did tech compare to energy 2017-2021?"

  # Resume an interrupted run (pass the run_id shown in the previous run's output)
  python portfolio_risk_workflow.py --run-id 20260427_143022

  # Resume automatically — re-run with the same query and the latest matching checkpoint is found
  python portfolio_risk_workflow.py --query "How did tech compare to energy 2017-2021?"

  # Delete checkpoint for a query and restart from scratch
  python portfolio_risk_workflow.py --query "How did tech compare to energy 2017-2021?" --reset

  # Delete a specific run by ID and restart
  python portfolio_risk_workflow.py --run-id 20260427_143022 --reset
        """,
    )
    parser.add_argument(
        "--query",
        type=str,
        default=DEFAULT_QUERY,
        help="Natural language research question about the S&P 500 (2017-2021 data)",
    )
    parser.add_argument(
        "--run-id",
        type=str,
        default=None,
        dest="run_id",
        help=(
            "Resume a specific interrupted run by its ID (the timestamp in the checkpoint "
            "filename, e.g. 20260427_143022). If omitted, the latest checkpoint matching "
            "the query is found automatically."
        ),
    )
    parser.add_argument(
        "--reset",
        action="store_true",
        default=False,
        help=(
            "Delete the checkpoint (and PPTX) for the matching query or run-id, "
            "then start a fresh run from Agent 1."
        ),
    )
    args = parser.parse_args()

    if args.reset:
        _delete_run(run_id=args.run_id, query=args.query if args.query != DEFAULT_QUERY or not args.run_id else None)

    executor = SP500AnalysisWorkflowExecutor()
    result = await executor.execute(args.query, run_id=None if args.reset else args.run_id)

    _banner("OUTPUT PATHS")
    print(f"PPTX Presentation -> {result['pptx_path'] or 'Not generated'}")
    print(f"Output JSON       -> {result['json_path']}")

    print("\n--- SLIDE DECK SPEC PREVIEW (first 600 chars) ---")
    print(result["slide_deck_spec"][:600])


if __name__ == "__main__":
    asyncio.run(main())
